#!/usr/bin/env python3
"""Build the Office-format data room (docx/pptx/xlsx) from the markdown sources.

Usage: python3 dataroom/scripts/build_office.py [output_dir]
Outputs one .docx per memo, a .pptx for the pitch deck, and a .xlsx for the
5-year financial model.
"""
import re
import sys
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.shared import Pt, RGBColor, Inches
from pptx import Presentation
from pptx.dml.color import RGBColor as PPTColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PInches, Pt as PPt
from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter

ROOT = Path(__file__).resolve().parent.parent
OUT = Path(sys.argv[1]) if len(sys.argv) > 1 else ROOT / "office"
OUT.mkdir(parents=True, exist_ok=True)

GOLD = "C9A227"
WARM_BLACK = "17120C"
INK = "231C13"

# ---------------------------------------------------------------- markdown --

INLINE_RE = re.compile(r"(\*\*.+?\*\*|\*[^*]+?\*|`[^`]+?`|\[[^\]]+?\]\([^)]+?\))")


def add_runs(paragraph, text, base_bold=False, base_italic=False):
    """Render markdown inline formatting (**bold**, *italic*, `code`, links) as runs."""
    for token in INLINE_RE.split(text):
        if not token:
            continue
        bold, italic, mono = base_bold, base_italic, False
        if token.startswith("**") and token.endswith("**") and len(token) > 4:
            token, bold = token[2:-2], True
        elif token.startswith("*") and token.endswith("*") and len(token) > 2:
            token, italic = token[1:-1], True
        elif token.startswith("`") and token.endswith("`"):
            token, mono = token[1:-1], True
        elif token.startswith("[") and "](" in token:
            m = re.match(r"\[([^\]]+)\]\(([^)]+)\)", token)
            if m:
                token = m.group(1)
        run = paragraph.add_run(token)
        run.bold = bold
        run.italic = italic
        if mono:
            run.font.name = "Consolas"
            run.font.size = Pt(9.5)


def strip_inline(text):
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*([^*]+?)\*", r"\1", text)
    text = re.sub(r"`([^`]+?)`", r"\1", text)
    text = re.sub(r"\[([^\]]+?)\]\([^)]+?\)", r"\1", text)
    return text


def parse_blocks(md_text):
    """Split markdown into blocks: heading/para/bullet/number/quote/code/table/hr."""
    lines = md_text.split("\n")
    blocks, i = [], 0
    while i < len(lines):
        line = lines[i]
        if line.startswith("```"):
            code, i = [], i + 1
            while i < len(lines) and not lines[i].startswith("```"):
                code.append(lines[i])
                i += 1
            i += 1
            blocks.append(("code", "\n".join(code)))
        elif line.strip().startswith("|") and i + 1 < len(lines) and set(lines[i + 1].replace("|", "").replace(":", "").strip()) <= {"-", " "} and lines[i + 1].strip().startswith("|"):
            rows = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                cells = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                rows.append(cells)
                i += 1
            del rows[1]  # separator row
            blocks.append(("table", rows))
        elif re.match(r"^#{1,6} ", line):
            level = len(line) - len(line.lstrip("#"))
            blocks.append(("heading", (level, line[level + 1:].strip())))
            i += 1
        elif line.strip() == "---":
            blocks.append(("hr", None))
            i += 1
        elif line.startswith("> "):
            quote = []
            while i < len(lines) and lines[i].startswith(">"):
                quote.append(lines[i][1:].strip())
                i += 1
            blocks.append(("quote", " ".join(q for q in quote if q)))
        elif re.match(r"^\s*[-*] ", line):
            m = re.match(r"^(\s*)[-*] (.*)$", line)
            blocks.append(("bullet", (len(m.group(1)) // 2, m.group(2))))
            i += 1
        elif re.match(r"^\s*\d+\. ", line):
            m = re.match(r"^\s*\d+\. (.*)$", line)
            blocks.append(("number", m.group(1)))
            i += 1
        elif line.strip():
            blocks.append(("para", line.strip()))
            i += 1
        else:
            i += 1
    return blocks


# ------------------------------------------------------------------- docx --

def style_docx(doc):
    normal = doc.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(10.5)
    for lvl, size in (("Title", 24), ("Heading 1", 16), ("Heading 2", 13), ("Heading 3", 11.5)):
        st = doc.styles[lvl]
        st.font.name = "Calibri"
        st.font.size = Pt(size)
        st.font.bold = True
        st.font.color.rgb = RGBColor.from_string(INK if lvl == "Title" else WARM_BLACK)


def set_cell_bg(cell, hex_color):
    shd = cell._tc.get_or_add_tcPr().makeelement(qn("w:shd"), {qn("w:val"): "clear", qn("w:fill"): hex_color})
    cell._tc.get_or_add_tcPr().append(shd)


def md_to_docx(md_path, out_path):
    doc = Document()
    style_docx(doc)
    for section in doc.sections:
        section.left_margin = section.right_margin = Inches(0.9)
    first_heading = True
    for kind, payload in parse_blocks(md_path.read_text()):
        if kind == "heading":
            level, text = payload
            text = strip_inline(text)
            if level == 1 and first_heading:
                doc.add_heading(text, 0)
                first_heading = False
            else:
                doc.add_heading(text, min(level, 4) - (1 if level > 1 else 0))
        elif kind == "para":
            p = doc.add_paragraph()
            add_runs(p, payload)
        elif kind == "quote":
            p = doc.add_paragraph(style="Intense Quote")
            add_runs(p, payload)
        elif kind == "bullet":
            indent, text = payload
            p = doc.add_paragraph(style="List Bullet 2" if indent else "List Bullet")
            add_runs(p, text)
        elif kind == "number":
            p = doc.add_paragraph(style="List Number")
            add_runs(p, payload)
        elif kind == "code":
            for code_line in payload.split("\n"):
                p = doc.add_paragraph()
                p.paragraph_format.space_after = Pt(0)
                run = p.add_run(code_line if code_line else " ")
                run.font.name = "Consolas"
                run.font.size = Pt(8)
        elif kind == "table":
            rows = payload
            ncols = max(len(r) for r in rows)
            table = doc.add_table(rows=len(rows), cols=ncols)
            table.style = "Table Grid"
            for ri, row in enumerate(rows):
                for ci in range(ncols):
                    cell = table.cell(ri, ci)
                    cell.paragraphs[0].text = ""
                    add_runs(cell.paragraphs[0], row[ci] if ci < len(row) else "", base_bold=(ri == 0))
                    for para in cell.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(9)
                    if ri == 0:
                        set_cell_bg(cell, "F2E8CF")
            doc.add_paragraph()
        elif kind == "hr":
            pass
    doc.save(out_path)


# ------------------------------------------------------------------- pptx --

def add_slide_bg(slide, hex_color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = PPTColor.from_string(hex_color)


def deck_to_pptx(md_path, out_path):
    text = md_path.read_text()
    slide_re = re.compile(r"^## Slide \d+ — (.+)$", re.M)
    titles = slide_re.findall(text)
    bodies = slide_re.split(text)[2::2] if slide_re.split(text) else []
    # slide_re.split gives [pre, title1, body1, title2, body2, ...]
    parts = slide_re.split(text)
    slides = list(zip(parts[1::2], parts[2::2]))

    prs = Presentation()
    prs.slide_width, prs.slide_height = PInches(13.333), PInches(7.5)
    blank = prs.slide_layouts[6]

    for idx, (title, body) in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        add_slide_bg(slide, WARM_BLACK)
        is_title_slide = idx == 0 or idx == len(slides) - 1

        # gold rule + title
        tbox = slide.shapes.add_textbox(PInches(0.6), PInches(0.35), PInches(12.1), PInches(1.0))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = strip_inline(title)
        run.font.size = PPt(40 if is_title_slide else 28)
        run.font.bold = True
        run.font.color.rgb = PPTColor.from_string(GOLD)
        run.font.name = "Georgia"
        if is_title_slide:
            p.alignment = PP_ALIGN.CENTER
            tbox.top = PInches(1.6)

        # body
        top = PInches(3.0 if is_title_slide else 1.45)
        box = slide.shapes.add_textbox(PInches(0.7), top, PInches(11.9), PInches(7.5 - (3.2 if is_title_slide else 1.6)))
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        table_blocks = []
        for kind, payload in parse_blocks(body):
            if kind == "table":
                table_blocks.append(payload)
                continue
            if kind in ("hr", "code"):
                if kind == "code":
                    for code_line in payload.split("\n"):
                        para = tf.paragraphs[0] if first else tf.add_paragraph()
                        first = False
                        run = para.add_run()
                        run.text = code_line or " "
                        run.font.name = "Consolas"
                        run.font.size = PPt(12)
                        run.font.color.rgb = PPTColor.from_string("E8DCC3")
                continue
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            if kind == "heading":
                _, htext = payload
                payload = htext
            if kind == "bullet":
                indent, btext = payload
                para.level = min(indent, 2)
                payload = "•  " + btext if indent == 0 else "–  " + btext
            if kind == "number":
                payload = payload  # keep text, numbered manually below
            clean = strip_inline(payload if isinstance(payload, str) else str(payload))
            run = para.add_run()
            run.text = clean
            emphasized = isinstance(payload, str) and payload.startswith("**")
            run.font.size = PPt(20 if (is_title_slide or emphasized) else 15)
            run.font.bold = emphasized
            run.font.color.rgb = PPTColor.from_string("F5EFDE" if not emphasized else "FFFFFF")
            run.font.name = "Calibri"
            if kind == "quote":
                run.font.italic = True
                run.font.color.rgb = PPTColor.from_string("D8C89A")
            para.space_after = PPt(8)
            if is_title_slide:
                para.alignment = PP_ALIGN.CENTER

        # tables (max 1 rendered per slide for space)
        for tdata in table_blocks[:1]:
            nrows, ncols = len(tdata), max(len(r) for r in tdata)
            theight = PInches(0.32) * nrows
            gshape = slide.shapes.add_table(nrows, ncols, PInches(0.7), PInches(7.3) - theight, PInches(11.9), theight)
            table = gshape.table
            for ri, row in enumerate(tdata):
                for ci in range(ncols):
                    cell = table.cell(ri, ci)
                    cell.text = strip_inline(row[ci]) if ci < len(row) else ""
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = PPt(11)
                            run.font.bold = ri == 0
                            run.font.color.rgb = PPTColor.from_string(WARM_BLACK if ri == 0 else "F5EFDE")
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = PPTColor.from_string(GOLD if ri == 0 else ("2A2318" if ri % 2 else "231C13"))
    prs.save(out_path)


# ------------------------------------------------------------------- xlsx --

HEAD_FILL = PatternFill("solid", fgColor=GOLD)
ALT_FILL = PatternFill("solid", fgColor="F7F1DF")
BOLD = Font(bold=True)
HEAD_FONT = Font(bold=True, color=WARM_BLACK)
THIN = Border(*[Side(style="thin", color="D9CFAF")] * 4)
YEARS = ["FY1", "FY2", "FY3", "FY4", "FY5"]


def sheet_header(ws, title, subtitle=None):
    ws["A1"] = title
    ws["A1"].font = Font(bold=True, size=14)
    if subtitle:
        ws["A2"] = subtitle
        ws["A2"].font = Font(italic=True, size=9, color="7A6A45")


def write_row(ws, r, label, values, fmt=None, bold=False, fill=None):
    ws.cell(row=r, column=1, value=label).font = BOLD if bold else Font()
    for i, v in enumerate(values):
        c = ws.cell(row=r, column=2 + i, value=v)
        if fmt:
            c.number_format = fmt
        if bold:
            c.font = BOLD
        if fill:
            c.fill = fill
        c.border = THIN
    ws.cell(row=r, column=1).border = THIN
    if fill:
        ws.cell(row=r, column=1).fill = fill


def year_header(ws, r):
    write_row(ws, r, "", YEARS, bold=True, fill=HEAD_FILL)
    for col in range(1, 7):
        ws.cell(row=r, column=col).font = HEAD_FONT


def build_xlsx(out_path):
    wb = Workbook()

    # -- Read Me
    ws = wb.active
    ws.title = "Read Me"
    sheet_header(ws, "FundExecs Technologies — 5-Year Financial Model", "Confidential · July 2026")
    notes = [
        "",
        "This workbook is an illustrative planning model, not a forecast. All figures are",
        "management estimates derived from the assumptions on the 'Assumptions' sheet and",
        "should be stress-tested in diligence. FY1 begins at close of the seed round",
        "(assumed Q4 2026). Figures in USD.",
        "",
        "Sheets:",
        "  Assumptions — pricing, mix, usage uplift, customers, churn, NRR, cost ratios",
        "  Revenue Build — blended ARPA and ARR computed by formula from Assumptions",
        "  P&L — recognized revenue, COGS, OpEx, EBITDA (formula-driven)",
        "  Headcount — hiring plan by function",
        "  Sensitivities — downside scenarios and mitigants",
        "",
        "Change any blue input cell on Assumptions and the model recalculates.",
    ]
    for i, line in enumerate(notes, start=3):
        ws.cell(row=i, column=1, value=line)
    ws.column_dimensions["A"].width = 100

    input_font = Font(color="1F4E79")

    # -- Assumptions
    ws = wb.create_sheet("Assumptions")
    sheet_header(ws, "Assumptions", "Blue cells are inputs — edit these; everything else is formula-driven.")
    r = 4
    year_header(ws, r); r += 1
    ws.cell(row=r, column=1, value="Pricing ($/org/month)").font = BOLD; r += 1
    price_rows = {}
    for tier, vals in (("Starter-tier equivalent", [50, 100, 125, 150, 175]),
                       ("Pro-tier equivalent", [300, 450, 600, 700, 800]),
                       ("Scale-tier equivalent", [1000, 1250, 1500, 1750, 2000])):
        write_row(ws, r, tier, vals, fmt="$#,##0")
        price_rows[tier.split("-")[0]] = r
        r += 1
    ws.cell(row=r, column=1, value="Customer mix (must sum to 100%)").font = BOLD; r += 1
    mix_rows = {}
    for tier, vals in (("Starter mix", [0.55, 0.50, 0.45, 0.40, 0.35]),
                       ("Pro mix", [0.35, 0.38, 0.40, 0.42, 0.43]),
                       ("Scale mix", [0.10, 0.12, 0.15, 0.18, 0.22])):
        write_row(ws, r, tier, vals, fmt="0%")
        mix_rows[tier.split(" ")[0]] = r
        r += 1
    ws.cell(row=r, column=1, value="Usage & customers").font = BOLD; r += 1
    uplift_row = r; write_row(ws, r, "Credits/usage uplift on subscription", [0.15, 0.20, 0.25, 0.30, 0.30], fmt="0%"); r += 1
    orgs_row = r; write_row(ws, r, "Paying orgs (EOY)", [40, 160, 450, 1000, 1900], fmt="#,##0"); r += 1
    churn_row = r; write_row(ws, r, "Gross logo churn (annual)", [0.15, 0.15, 0.12, 0.10, 0.10], fmt="0%"); r += 1
    nrr_row = r; write_row(ws, r, "Net revenue retention", [1.00, 1.10, 1.15, 1.18, 1.20], fmt="0%"); r += 1
    ws.cell(row=r, column=1, value="Other revenue & cost ratios").font = BOLD; r += 1
    mkt_row = r; write_row(ws, r, "Marketplace + data/API ARR ($M, EOY)", [0, 0.05, 0.55, 1.5, 3.0], fmt="$0.00"); r += 1
    cogs_row = r; write_row(ws, r, "COGS as % of revenue", [0.55, 0.31, 0.23, 0.21, 0.19], fmt="0%"); r += 1
    ws.cell(row=r, column=1, value="OpEx plan ($M)").font = BOLD; r += 1
    rd_row = r; write_row(ws, r, "R&D / Engineering", [0.9, 1.5, 2.6, 4.5, 7.0], fmt="$0.00"); r += 1
    sm_row = r; write_row(ws, r, "Sales & Marketing", [0.2, 0.5, 1.3, 2.7, 4.9], fmt="$0.00"); r += 1
    ga_row = r; write_row(ws, r, "G&A (incl. legal/compliance)", [0.3, 0.5, 0.9, 1.6, 2.5], fmt="$0.00"); r += 1
    for row_idx in [price_rows["Starter"], price_rows["Pro"], price_rows["Scale"],
                    mix_rows["Starter"], mix_rows["Pro"], mix_rows["Scale"],
                    uplift_row, orgs_row, churn_row, nrr_row, mkt_row, cogs_row, rd_row, sm_row, ga_row]:
        for col in range(2, 7):
            ws.cell(row=row_idx, column=col).font = input_font
    ws.column_dimensions["A"].width = 42
    for col in range(2, 7):
        ws.column_dimensions[get_column_letter(col)].width = 12

    A = "Assumptions"

    # -- Revenue Build
    ws = wb.create_sheet("Revenue Build")
    sheet_header(ws, "Revenue Build", "Formula-driven from the Assumptions sheet.")
    r = 4
    year_header(ws, r); r += 1
    arpa_row = r
    for i in range(5):
        col = get_column_letter(2 + i)
        formula = (f"=({A}!{col}{price_rows['Starter']}*{A}!{col}{mix_rows['Starter']}"
                   f"+{A}!{col}{price_rows['Pro']}*{A}!{col}{mix_rows['Pro']}"
                   f"+{A}!{col}{price_rows['Scale']}*{A}!{col}{mix_rows['Scale']})*12*(1+{A}!{col}{uplift_row})")
        c = ws.cell(row=r, column=2 + i, value=formula)
        c.number_format = "$#,##0"
        c.border = THIN
    ws.cell(row=r, column=1, value="Blended ARPA ($/org/yr)").font = BOLD
    ws.cell(row=r, column=1).border = THIN
    r += 1
    sub_row = r
    for i in range(5):
        col = get_column_letter(2 + i)
        c = ws.cell(row=r, column=2 + i, value=f"={A}!{col}{orgs_row}*{get_column_letter(2+i)}{arpa_row}/1000000")
        c.number_format = "$0.00"
        c.border = THIN
    ws.cell(row=r, column=1, value="Subscription + usage ARR ($M, EOY)").border = THIN
    r += 1
    mkt2_row = r
    for i in range(5):
        col = get_column_letter(2 + i)
        c = ws.cell(row=r, column=2 + i, value=f"={A}!{col}{mkt_row}")
        c.number_format = "$0.00"
        c.border = THIN
    ws.cell(row=r, column=1, value="Marketplace + data/API ARR ($M)").border = THIN
    r += 1
    total_arr_row = r
    write_row(ws, r, "Total ARR ($M, EOY)", [f"={get_column_letter(2+i)}{sub_row}+{get_column_letter(2+i)}{mkt2_row}" for i in range(5)], fmt="$0.00", bold=True, fill=ALT_FILL)
    r += 1
    rec_rev_row = r
    vals = []
    for i in range(5):
        col = get_column_letter(2 + i)
        prev = f"{get_column_letter(1+i)}{total_arr_row}" if i else "0"
        vals.append(f"=({prev}+{col}{total_arr_row})/2*0.95")
    write_row(ws, r, "Recognized revenue ($M) ≈ avg ARR × 95%", vals, fmt="$0.00", bold=True)
    ws.column_dimensions["A"].width = 42
    for col in range(2, 7):
        ws.column_dimensions[get_column_letter(col)].width = 12

    RB = "'Revenue Build'"

    # -- P&L
    ws = wb.create_sheet("P&L")
    sheet_header(ws, "P&L Summary ($M)", "Formula-driven; EBITDA breakeven expected late FY5 on plan.")
    r = 4
    year_header(ws, r); r += 1
    rev_row = r
    write_row(ws, r, "Revenue", [f"={RB}!{get_column_letter(2+i)}{rec_rev_row}" for i in range(5)], fmt="$0.00"); r += 1
    cogs2_row = r
    write_row(ws, r, "COGS (AI compute, infra, support)", [f"={get_column_letter(2+i)}{rev_row}*{A}!{get_column_letter(2+i)}{cogs_row}" for i in range(5)], fmt="$0.00"); r += 1
    gp_row = r
    write_row(ws, r, "Gross profit", [f"={get_column_letter(2+i)}{rev_row}-{get_column_letter(2+i)}{cogs2_row}" for i in range(5)], fmt="$0.00", bold=True, fill=ALT_FILL); r += 1
    write_row(ws, r, "Gross margin", [f"=IF({get_column_letter(2+i)}{rev_row}=0,0,{get_column_letter(2+i)}{gp_row}/{get_column_letter(2+i)}{rev_row})" for i in range(5)], fmt="0%"); r += 1
    rd2 = r; write_row(ws, r, "R&D / Engineering", [f"={A}!{get_column_letter(2+i)}{rd_row}" for i in range(5)], fmt="$0.00"); r += 1
    sm2 = r; write_row(ws, r, "Sales & Marketing", [f"={A}!{get_column_letter(2+i)}{sm_row}" for i in range(5)], fmt="$0.00"); r += 1
    ga2 = r; write_row(ws, r, "G&A", [f"={A}!{get_column_letter(2+i)}{ga_row}" for i in range(5)], fmt="$0.00"); r += 1
    opex_row = r
    write_row(ws, r, "Total OpEx", [f"=SUM({get_column_letter(2+i)}{rd2}:{get_column_letter(2+i)}{ga2})" for i in range(5)], fmt="$0.00", bold=True); r += 1
    write_row(ws, r, "EBITDA", [f"={get_column_letter(2+i)}{gp_row}-{get_column_letter(2+i)}{opex_row}" for i in range(5)], fmt="$0.00;($0.00)", bold=True, fill=ALT_FILL)
    ws.column_dimensions["A"].width = 42
    for col in range(2, 7):
        ws.column_dimensions[get_column_letter(col)].width = 12

    # -- Headcount
    ws = wb.create_sheet("Headcount")
    sheet_header(ws, "Headcount Plan (FTEs, EOY)")
    r = 4
    year_header(ws, r); r += 1
    eng = r; write_row(ws, r, "Engineering & product", [4, 7, 12, 20, 30], fmt="#,##0"); r += 1
    gtm = r; write_row(ws, r, "GTM & success", [1, 3, 7, 13, 21], fmt="#,##0"); r += 1
    ops = r; write_row(ws, r, "G&A / ops", [1, 2, 3, 5, 8], fmt="#,##0"); r += 1
    write_row(ws, r, "Total", [f"=SUM({get_column_letter(2+i)}{eng}:{get_column_letter(2+i)}{ops})" for i in range(5)], fmt="#,##0", bold=True, fill=ALT_FILL)
    ws.column_dimensions["A"].width = 42
    for col in range(2, 7):
        ws.column_dimensions[get_column_letter(col)].width = 12

    # -- Sensitivities
    ws = wb.create_sheet("Sensitivities")
    sheet_header(ws, "Key Sensitivities")
    rows = [
        ("Variable", "Downside effect", "Mitigant"),
        ("Repricing lags proof of value", "FY3+ ARR −40–50%", "Credits still capture usage growth; delay hires to match"),
        ("Churn 20%+ in beachhead", "FY5 orgs ~1,300", "Compounding graphs/Brains raise switching costs with tenure"),
        ("AI compute cost doesn't fall", "Gross margin plateaus ~70%", "Cost-tiered routing (live), per-run caps, pass-through via credits"),
        ("Slower logo adds (half pace)", "FY5 ARR ~$14M", "Still a strong Series B profile; burn scales down with GTM spend"),
    ]
    for ri, row in enumerate(rows, start=4):
        for ci, val in enumerate(row, start=1):
            c = ws.cell(row=ri, column=ci, value=val)
            c.border = THIN
            c.alignment = Alignment(wrap_text=True, vertical="top")
            if ri == 4:
                c.font = HEAD_FONT
                c.fill = HEAD_FILL
    ws.column_dimensions["A"].width = 32
    ws.column_dimensions["B"].width = 28
    ws.column_dimensions["C"].width = 60

    wb.save(out_path)


# ------------------------------------------------------------------- main --

DOCX_SOURCES = [
    "00-data-room-index.md", "01-executive-summary.md", "03-product-vision.md",
    "04-product-roadmap.md", "05-platform-architecture-map.md",
    "06-ai-executive-team-earn-agent-spec.md", "07-business-model-canvas.md",
    "09-use-of-funds.md", "10-market-opportunity-memo.md",
    "11-competitive-landscape.md", "12-gtm-strategy.md",
    "13-risk-compliance-memo.md", "14-bgi-fund-i-relationship-memo.md",
]

if __name__ == "__main__":
    for name in DOCX_SOURCES:
        src = ROOT / name
        dst = OUT / (src.stem + ".docx")
        md_to_docx(src, dst)
        print(f"docx  {dst.name}")
    deck_to_pptx(ROOT / "02-pitch-deck.md", OUT / "02-pitch-deck.pptx")
    print("pptx  02-pitch-deck.pptx")
    build_xlsx(OUT / "08-five-year-financial-model.xlsx")
    print("xlsx  08-five-year-financial-model.xlsx")
    print(f"\nDone → {OUT}")
